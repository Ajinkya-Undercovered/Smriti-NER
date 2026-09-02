import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sih_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Colors
    NAVY_BLUE = RGBColor(27, 54, 93)     # #1B365D (SIH Primary)
    SKY_BLUE = RGBColor(2, 132, 199)     # #0284C7
    ORANGE = RGBColor(234, 88, 12)       # #EA580C
    DARK_TEXT = RGBColor(30, 41, 59)     # #1E293B
    LIGHT_BG = RGBColor(248, 250, 252)   # #F8FAFC
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(100, 116, 139)  # #64748B
    CARD_BG = RGBColor(241, 245, 249)    # #F1F5F9

    image_path = r"C:\Users\ajink\.gemini\antigravity\brain\35e1dec1-288c-467c-830d-9f1f5ba63403\smriti_ner_clean_concept_1788257265433.jpg"

    def add_common_decorations(slide, slide_num, title_text):
        # Bottom Blue Accent Bar
        bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = SKY_BLUE
        bottom_bar.line.fill.background()
        
        # Bottom text
        tf = bottom_bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"@SIH Idea submission- Template {slide_num}"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Top Left Team Name Badge
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.35), Inches(1.8), Inches(0.8))
        oval.fill.solid()
        oval.fill.fore_color.rgb = WHITE
        oval.line.color.rgb = NAVY_BLUE
        oval.line.width = Pt(1.5)
        tf_o = oval.text_frame
        p_o = tf_o.paragraphs[0]
        p_o.text = "Your Team\nName"
        p_o.font.size = Pt(11)
        p_o.font.bold = True
        p_o.font.color.rgb = NAVY_BLUE
        p_o.alignment = PP_ALIGN.CENTER

        # Top Right SIH Header
        txBox = slide.shapes.add_textbox(Inches(10.2), Inches(0.3), Inches(2.6), Inches(0.9))
        tf_s = txBox.text_frame
        p_s1 = tf_s.paragraphs[0]
        p_s1.text = "SMART INDIA"
        p_s1.font.size = Pt(14)
        p_s1.font.bold = True
        p_s1.font.color.rgb = NAVY_BLUE
        p_s1.alignment = PP_ALIGN.RIGHT
        p_s2 = tf_s.add_paragraph()
        p_s2.text = "HACKATHON 2025"
        p_s2.font.size = Pt(14)
        p_s2.font.bold = True
        p_s2.font.color.rgb = ORANGE
        p_s2.alignment = PP_ALIGN.RIGHT

        # Center Title
        title_box = slide.shapes.add_textbox(Inches(2.6), Inches(0.35), Inches(7.4), Inches(0.8))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_BLUE
        p_t.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    
    # Title Header Box
    s1_header = s1.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
    tf1 = s1_header.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2025"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_BLUE
    p1.alignment = PP_ALIGN.CENTER
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "TITLE PAGE"
    p1_sub.font.size = Pt(20)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = ORANGE
    p1_sub.alignment = PP_ALIGN.CENTER

    # Details Card Container
    s1_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    s1_card.fill.solid()
    s1_card.fill.fore_color.rgb = CARD_BG
    s1_card.line.color.rgb = SKY_BLUE
    s1_card.line.width = Pt(1.5)

    tf_card = s1_card.text_frame
    tf_card.word_wrap = True

    entries = [
        ("Problem Statement ID :", " Enter your assigned PS ID (e.g. SIH2025-MED01)"),
        ("Problem Statement Title :", " AI-Driven Culturally Grounded Digital Therapeutics & Cognitive Care for Dementia"),
        ("Project Name :", " Smriti-NER (স্মৃতি) — Cognitive Health Platform for North East India"),
        ("Theme :", " MedTech / BioTech / Healthcare & Rural Health"),
        ("PS Category :", " Software"),
        ("Team ID :", " Enter your Team ID"),
        ("Team Name :", " Enter your Registered Team Name")
    ]

    for idx, (label, val) in enumerate(entries):
        p = tf_card.paragraphs[0] if idx == 0 else tf_card.add_paragraph()
        p.space_after = Pt(10)
        run_l = p.add_run()
        run_l.text = f"• {label}"
        run_l.font.bold = True
        run_l.font.size = Pt(15)
        run_l.font.color.rgb = NAVY_BLUE

        run_v = p.add_run()
        run_v.text = val
        run_v.font.size = Pt(15)
        run_v.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_common_decorations(s2, 2, "IDEA TITLE: SMRITI-NER (স্মৃতি)")

    # Left Column: Structured Solution & Value Proposition (Width: 6.8 Inches)
    s2_box = s2.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(6.8), Inches(5.5))
    tf2 = s2_box.text_frame
    tf2.word_wrap = True

    # 1. Proposed Solution
    p = tf2.paragraphs[0]
    p.text = "Proposed Solution (3 Core Pillars):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    bullets_s2_1 = [
        ("Zero-Confusion Elderly UI: ", "5 giant high-affordance tactile tiles with password-free 1-tap fast-pass login for dementia patients."),
        ("8-Game Cultural Therapeutic Suite: ", "Clinically designed cognitive exercises using familiar North Eastern motifs (Kaziranga rhino, tea sorting, Bihu rhythms) stimulating memory and combating visual agnosia."),
        ("Multilingual Neural AI Voice Companion: ", "Empathetic voice assistant providing hands-free guidance and daily medication, water, and appointment reminders in Assamese and English.")
    ]
    for b_title, b_desc in bullets_s2_1:
        p = tf2.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {b_title}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = ORANGE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_TEXT

    # 2. How it addresses the problem
    p = tf2.add_paragraph()
    p.space_before = Pt(6)
    p.text = "How It Addresses The Problem:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    p_p = tf2.add_paragraph()
    p_p.space_after = Pt(4)
    r1 = p_p.add_run()
    r1.text = "• Bridges Rural Diagnostic Gap: "
    r1.font.bold = True
    r1.font.size = Pt(11)
    r1.font.color.rgb = SKY_BLUE
    r2 = p_p.add_run()
    r2.text = "Solves the 85%+ undiagnosed dementia gap in rural NER with zero-cost smartphone deployment and empowers local ASHA health workers."
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK_TEXT

    # 3. Innovation & Uniqueness
    p = tf2.add_paragraph()
    p.space_before = Pt(6)
    p.text = "Innovation & Uniqueness:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    p_i = tf2.add_paragraph()
    r1 = p_i.add_run()
    r1.text = "• Passive Micro-Latency Tracking: "
    r1.font.bold = True
    r1.font.size = Pt(11)
    r1.font.color.rgb = ORANGE
    r2 = p_i.add_run()
    r2.text = "Records millisecond reaction speeds and accuracy curves to generate MMSE-correlated clinical radar reports for neurologists."
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK_TEXT

    # Right Column: Professional Graphic
    if os.path.exists(image_path):
        s2.shapes.add_picture(image_path, Inches(7.6), Inches(1.5), width=Inches(5.2))

    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_common_decorations(s3, 3, "TECHNICAL APPROACH")

    # Box 1: Tech Stack
    s3_t1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.3))
    s3_t1.fill.solid()
    s3_t1.fill.fore_color.rgb = CARD_BG
    s3_t1.line.color.rgb = SKY_BLUE
    tf_t1 = s3_t1.text_frame
    tf_t1.word_wrap = True
    
    p = tf_t1.paragraphs[0]
    p.text = "Technologies & Architecture:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(10)

    techs = [
        ("Frontend / UX: ", "React 19, Vite, TailwindCSS (WCAG 2.1 AAA High-Contrast Palette for Elderly)."),
        ("Voice AI Engine: ", "Multilingual Neural Speech Synthesis, Web Speech API, Web Audio API."),
        ("Backend & Database: ", "Supabase Cloud (PostgreSQL), Granular Row-Level Security (RLS)."),
        ("Clinical Analytics: ", "Recharts (5-axis Cognitive Domain Radar Charts, Longitudinal Latency Curves)."),
        ("Deployment: ", "Vercel Edge Network, Progressive Web App (PWA) with Offline IndexedDB Cache.")
    ]
    for k, v in techs:
        p = tf_t1.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = f"• {k}"
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = ORANGE
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_TEXT

    # Box 2: Methodology & Workflow
    s3_t2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.3))
    s3_t2.fill.solid()
    s3_t2.fill.fore_color.rgb = CARD_BG
    s3_t2.line.color.rgb = SKY_BLUE
    tf_t2 = s3_t2.text_frame
    tf_t2.word_wrap = True

    p = tf_t2.paragraphs[0]
    p.text = "Implementation Workflow & Process:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(10)

    steps = [
        ("1. 1-Tap Access: ", "Dementia patient logs in with single-tap remembered avatar, zero password friction."),
        ("2. Multilingual Engagement: ", "Listens to natural Assamese/English audio orientation and selects from 8 therapeutic games."),
        ("3. Passive Metric Extraction: ", "Engine logs interaction latencies (ms), decision errors, and motor accuracy in real-time."),
        ("4. Smart Reminders: ", "Proactive automated voice reminders for daily medications, hydration, and doctor visits."),
        ("5. Clinical Tele-Triaging: ", "Syncs encrypted longitudinal records to ASHA workers and GMCH neurologists for early intervention.")
    ]
    for k, v in steps:
        p = tf_t2.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = f"• {k}"
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = SKY_BLUE
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_common_decorations(s4, 4, "FEASIBILITY AND VIABILITY")

    s4_box = s4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf4 = s4_box.text_frame
    tf4.word_wrap = True

    # 1. Feasibility
    p = tf4.paragraphs[0]
    p.text = "1. Feasibility Analysis of the Idea:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(4)

    feas_points = [
        "• Working Live Prototype: Already fully developed and deployed on Vercel with real-time cloud data sync.",
        "• Zero Hardware Cost: Runs seamlessly on any standard smartphone, tablet, or PHC desktop without expensive external sensors.",
        "• Low-Bandwidth Optimization: Designed to operate reliably on rural 2G/3G networks across remote hilly regions."
    ]
    for pt in feas_points:
        p = tf4.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(3)

    # 2. Potential Challenges & Strategies (Table Layout via Text)
    p = tf4.add_paragraph()
    p.space_before = Pt(8)
    p.text = "2. Potential Challenges & Mitigation Strategies:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(6)

    risks = [
        ("Risk 1: Intermittent Internet in Remote Hilly Districts (e.g., Majuli, Dima Hasao)", 
         "Mitigation: Offline-First PWA architecture with local IndexedDB storage that automatically syncs when online."),
        ("Risk 2: Digital Hesitancy & Tech Anxiety among 70+ Dementia Seniors", 
         "Mitigation: High-affordance 5-tile UI, giant touch targets, voice guidance, and zero password typing."),
        ("Risk 3: Health Data Privacy & Security (DPDP Act / Healthcare Compliance)", 
         "Mitigation: Granular Row-Level Security (RLS), encrypted patient IDs, and anonymized clinical session logs.")
    ]
    for r, m in risks:
        p = tf4.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"• {r}\n  ➔ "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = ORANGE
        r2 = p.add_run()
        r2.text = m
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_common_decorations(s5, 5, "IMPACT AND BENEFITS")

    # 3 Cards Layout
    cards_data = [
        ("Target Audience Impact", SKY_BLUE, [
            "Slowing Cognitive Decline: Regular cognitive stimulation preserves neural plasticity in MCI patients.",
            "Emotional Calming & Sundowning Relief: Traditional folk music and bird calls reduce anxiety and agitation.",
            "Independence & Dignity: Seniors independently track water, medicines, and daily orientation."
        ]),
        ("Social & Healthcare Benefits", ORANGE, [
            "Empowering ASHA Workers: Provides objective digital screening tools during routine rural home visits.",
            "Bridging Healthcare Divide: Direct digital linkage from village PHCs to tertiary centers (GMCH/NEIGRIHMS).",
            "Culturally Inclusive Care: First healthcare tool embracing indigenous North Eastern heritage and languages."
        ]),
        ("Economic Viability", NAVY_BLUE, [
            "Massive Cost Savings: Reduces avoidable hospitalizations and specialized dementia caregiver costs.",
            "National Scalability: Seamlessly aligns with National Programme for Health Care of the Elderly (NPHCE).",
            "Zero Cost to Seniors: High public health ROI with zero proprietary hardware requirements."
        ])
    ]

    for idx, (c_title, c_color, c_bullets) in enumerate(cards_data):
        left_pos = Inches(0.6 + idx * 4.1)
        card_shp = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.4), Inches(3.9), Inches(5.3))
        card_shp.fill.solid()
        card_shp.fill.fore_color.rgb = CARD_BG
        card_shp.line.color.rgb = c_color
        card_shp.line.width = Pt(2)
        
        tf_c = card_shp.text_frame
        tf_c.word_wrap = True
        p = tf_c.paragraphs[0]
        p.text = c_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = c_color
        p.space_after = Pt(12)

        for b in c_bullets:
            p = tf_c.add_paragraph()
            p.space_after = Pt(8)
            p.text = f"• {b}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_common_decorations(s6, 6, "RESEARCH AND REFERENCES")

    s6_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf6 = s6_box.text_frame
    tf6.word_wrap = True

    p = tf6.paragraphs[0]
    p.text = "Clinical Research & Scientific Grounding:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(4)

    refs = [
        ("World Health Organization (WHO): ", "Global Action Plan on the Public Health Response to Dementia (2017–2025)."),
        ("Dementia India Report (ARDSI): ", "Assessment of high prevalence and clinical diagnostic gaps in rural and North Eastern populations."),
        ("The Lancet Commission (2024): ", "Dementia prevention, intervention, and care — emphasizing early cognitive and sensory stimulation."),
        ("Folstein et al. (MMSE & MoCA): ", "Mini-Mental State Examination protocols and adaptation for multilingual Indian demographics.")
    ]
    for k, v in refs:
        p = tf6.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {k}"
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = ORANGE
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_TEXT

    p = tf6.add_paragraph()
    p.space_before = Pt(8)
    p.text = "Standards & Frameworks:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(4)

    stds = [
        ("W3C Web Content Accessibility Guidelines (WCAG 2.1 AAA): ", "Elderly touch targets, color contrast, and audio guidance standards."),
        ("Ayushman Bharat Digital Mission (ABDM): ", "Healthcare registry interoperability and tele-monitoring frameworks.")
    ]
    for k, v in stds:
        p = tf6.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {k}"
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = SKY_BLUE
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_TEXT

    p = tf6.add_paragraph()
    p.space_before = Pt(8)
    p.text = "Live Prototype & Verification Link:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.space_after = Pt(4)

    p_url = tf6.add_paragraph()
    r1 = p_url.add_run()
    r1.text = "• Live Application URL: "
    r1.font.bold = True
    r1.font.size = Pt(12)
    r1.font.color.rgb = NAVY_BLUE
    r2 = p_url.add_run()
    r2.text = "https://smriti-ner-three.vercel.app (Verified Working Production Prototype)"
    r2.font.size = Pt(12)
    r2.font.bold = True
    r2.font.color.rgb = SKY_BLUE

    # Save output file
    output_filename = "Smriti_NER_SIH_2025_Submission.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'")

if __name__ == "__main__":
    create_sih_presentation()
